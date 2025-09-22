# import os
# import magic
# from typing import Optional, Tuple
# from langchain.document_loaders import PyPDFLoader, TextLoader
# from langchain.document_loaders.word_document import UnstructuredWordDocumentLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from PIL import Image
# import logging

# logger = logging.getLogger(__name__)

# class FileProcessor:
#     def __init__(self):
#         self.text_splitter = RecursiveCharacterTextSplitter(
#             chunk_size=1000,
#             chunk_overlap=200,
#             length_function=len,
#         )
    
#     async def process_file(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process uploaded file and extract content"""
#         try:
#             file_type = magic.from_file(file_path, mime=True)
            
#             if file_type.startswith('image/'):
#                 return await self._process_image(file_path, filename)
#             elif file_type == 'application/pdf':
#                 return await self._process_pdf(file_path, filename)
#             elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
#                               'application/msword']:
#                 return await self._process_word(file_path, filename)
#             elif file_type.startswith('text/'):
#                 return await self._process_text(file_path, filename)
#             else:
#                 return f"Unsupported file type: {file_type}", "error"
                
#         except Exception as e:
#             logger.error(f"Error processing file {filename}: {e}")
#             return f"Error processing file: {str(e)}", "error"
    
#     async def _process_image(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process image file"""
#         try:
#             with Image.open(file_path) as img:
#                 width, height = img.size
#                 format_name = img.format
                
#             content = f"Image uploaded: {filename}\n"
#             content += f"Format: {format_name}\n"
#             content += f"Dimensions: {width}x{height} pixels\n"
#             content += "Ready for analysis. You can ask me to describe the image or answer questions about it."
            
#             return content, "image"
#         except Exception as e:
#             return f"Error processing image: {str(e)}", "error"
    
#     async def _process_pdf(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process PDF file"""
#         try:
#             loader = PyPDFLoader(file_path)
#             documents = loader.load()
            
#             if not documents:
#                 return "PDF appears to be empty or unreadable.", "error"
            
#             # Combine all pages
#             full_text = "\n\n".join([doc.page_content for doc in documents])
            
#             # Create summary
#             content = f"PDF Document: {filename}\n"
#             content += f"Pages: {len(documents)}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"
            
#             if len(full_text) > 2000:
#                 content += "Document preview (first 2000 characters):\n"
#                 content += full_text[:2000] + "...\n\n"
#             else:
#                 content += "Full document content:\n"
#                 content += full_text + "\n\n"
            
#             content += "Document loaded successfully. You can ask me questions about its content."
            
#             return content, "pdf"
#         except Exception as e:
#             return f"Error processing PDF: {str(e)}", "error"
    
#     async def _process_word(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process Word document"""
#         try:
#             loader = UnstructuredWordDocumentLoader(file_path)
#             documents = loader.load()
            
#             if not documents:
#                 return "Word document appears to be empty or unreadable.", "error"
            
#             full_text = "\n\n".join([doc.page_content for doc in documents])
            
#             content = f"Word Document: {filename}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"
            
#             if len(full_text) > 2000:
#                 content += "Document preview (first 2000 characters):\n"
#                 content += full_text[:2000] + "...\n\n"
#             else:
#                 content += "Full document content:\n"
#                 content += full_text + "\n\n"
            
#             content += "Document loaded successfully. You can ask me questions about its content."
            
#             return content, "word"
#         except Exception as e:
#             return f"Error processing Word document: {str(e)}", "error"
    
#     async def _process_text(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process text file"""
#         try:
#             loader = TextLoader(file_path, encoding='utf-8')
#             documents = loader.load()
            
#             if not documents:
#                 return "Text file appears to be empty.", "error"
            
#             full_text = documents[0].page_content
            
#             content = f"Text File: {filename}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"
            
#             if len(full_text) > 2000:
#                 content += "File preview (first 2000 characters):\n"
#                 content += full_text[:2000] + "...\n\n"
#             else:
#                 content += "Full file content:\n"
#                 content += full_text + "\n\n"
            
#             content += "File loaded successfully. You can ask me questions about its content."
            
#             return content, "text"
#         except Exception as e:
#             return f"Error processing text file: {str(e)}", "error"

# file_processor = FileProcessor()

# import os
# import magic
# import mimetypes  
# import pandas as pd
# from PIL import Image
# from typing import Optional, Tuple
# from pptx import Presentation
# from langchain.document_loaders import PyPDFLoader, TextLoader
# from langchain.document_loaders.word_document import UnstructuredWordDocumentLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# import logging

# logger = logging.getLogger(__name__)


# class FileProcessor:
#     def __init__(self):
#         self.text_splitter = RecursiveCharacterTextSplitter(
#             chunk_size=1000,
#             chunk_overlap=200,
#             length_function=len,
#         )

#     async def process_file(self, file_path: str, filename: str) -> Tuple[str, str]:
#         """Process uploaded file and extract content summary"""
#         try:
#             try:
#                file_type = magic.from_file(file_path, mime=True)
#             except Exception as e:
#                 logger.warning(f"[magic] failed to detect file type: {e}")
#                 file_type, _ = mimetypes.guess_type(file_path)
#                 if not file_type:
#                    file_type = "application/octet-stream"

#             logger.info(f"Detected file type for {filename}: {file_type}")

#             if file_type.startswith("image/"):
#                 return await self._process_image(file_path, filename)
#             elif file_type == "application/pdf":
#                 return await self._process_pdf(file_path, filename)
#             elif file_type in [
#                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
#                 "application/msword",
#             ]:
#                 return await self._process_word(file_path, filename)
#             elif file_type == "text/plain":
#                 return await self._process_text(file_path, filename)
#             elif file_type == "text/csv":
#                 return await self._process_csv(file_path, filename)
#             elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
#                 return await self._process_excel(file_path, filename)
#             elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
#                 return await self._process_ppt(file_path, filename)
#             elif file_type == "application/octet-stream":
#                 return await self._process_unknown(file_path, filename, file_type)
#             else:
#                 return f"Unsupported file type: {file_type}", "error"


#         except Exception as e:
#             logger.error(f"Error processing file {filename}: {e}")
#             return f"Error processing file: {str(e)}", "error"

#     async def _process_image(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             with Image.open(file_path) as img:
#                 width, height = img.size
#                 format_name = img.format

#             content = f"Image uploaded: {filename}\n"
#             content += f"Format: {format_name}\n"
#             content += f"Dimensions: {width}x{height} pixels\n"
#             content += "Ready for analysis. You can ask me to describe the image or answer questions about it."

#             return content, "image"
#         except Exception as e:
#             return f"Error processing image: {str(e)}", "error"

#     async def _process_pdf(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             loader = PyPDFLoader(file_path)
#             documents = loader.load()

#             if not documents:
#                 return "PDF appears to be empty or unreadable.", "error"

#             full_text = "\n\n".join([doc.page_content for doc in documents])

#             content = f"PDF Document: {filename}\n"
#             content += f"Pages: {len(documents)}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"

#             preview = full_text[:2000] + "..." if len(full_text) > 2000 else full_text
#             content += "Document preview:\n" + preview + "\n\n"
#             content += "Document loaded successfully. You can ask me questions about its content."

#             return content, "pdf"
#         except Exception as e:
#             return f"Error processing PDF: {str(e)}", "error"

#     async def _process_word(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             loader = UnstructuredWordDocumentLoader(file_path)
#             documents = loader.load()

#             if not documents:
#                 return "Word document appears to be empty or unreadable.", "error"

#             full_text = "\n\n".join([doc.page_content for doc in documents])

#             content = f"Word Document: {filename}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"
#             preview = full_text[:2000] + "..." if len(full_text) > 2000 else full_text
#             content += "Document preview:\n" + preview + "\n\n"
#             content += "Document loaded successfully. You can ask me questions about its content."

#             return content, "word"
#         except Exception as e:
#             return f"Error processing Word document: {str(e)}", "error"

#     async def _process_text(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             loader = TextLoader(file_path, encoding="utf-8")
#             documents = loader.load()

#             if not documents:
#                 return "Text file appears to be empty.", "error"

#             full_text = documents[0].page_content

#             content = f"Text File: {filename}\n"
#             content += f"Content length: {len(full_text)} characters\n\n"
#             preview = full_text[:2000] + "..." if len(full_text) > 2000 else full_text
#             content += "File preview:\n" + preview + "\n\n"
#             content += "File loaded successfully. You can ask me questions about its content."

#             return content, "text"
#         except Exception as e:
#             return f"Error processing text file: {str(e)}", "error"

#     async def _process_csv(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             df = pd.read_csv(file_path)
#             content = f"CSV File: {filename}\n"
#             content += f"Rows: {df.shape[0]}, Columns: {df.shape[1]}\n\n"
#             content += f"Preview:\n{df.head(5).to_string(index=False)}\n\n"
#             content += "CSV file loaded successfully. You can ask me about its contents."
#             return content, "csv"
#         except Exception as e:
#             return f"Error processing CSV file: {str(e)}", "error"

#     async def _process_excel(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             df = pd.read_excel(file_path)
#             content = f"Excel File: {filename}\n"
#             content += f"Rows: {df.shape[0]}, Columns: {df.shape[1]}\n\n"
#             content += f"Preview:\n{df.head(5).to_string(index=False)}\n\n"
#             content += "Excel file loaded successfully. You can ask me about its contents."
#             return content, "excel"
#         except Exception as e:
#             return f"Error processing Excel file: {str(e)}", "error"

#     async def _process_ppt(self, file_path: str, filename: str) -> Tuple[str, str]:
#         try:
#             prs = Presentation(file_path)
#             content = f"PowerPoint File: {filename}\n"
#             content += f"Slides: {len(prs.slides)}\n\n"
#             text = ""
#             for i, slide in enumerate(prs.slides[:3]):
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         text += f"Slide {i+1}: {shape.text}\n"
#             content += text[:2000] + "...\n\n" if len(text) > 2000 else text
#             content += "PowerPoint file loaded. Ask me about its content."
#             return content, "ppt"
#         except Exception as e:
#             return f"Error processing PowerPoint file: {str(e)}", "error"

#     async def _process_unknown(self, file_path: str, filename: str, file_type: str) -> Tuple[str, str]:
#         try:
#             size = os.path.getsize(file_path)
#             content = f"File uploaded: {filename}\n"
#             content += f"Type: {file_type}\n"
#             content += f"Size: {round(size / 1024, 2)} KB\n\n"
#             content += "This file type is not supported for detailed parsing. You can try converting it to PDF, TXT, or DOCX."
#             return content, "unsupported"
#         except Exception as e:
#             return f"Error analyzing unknown file type: {str(e)}", "error"


# # Export instance
# file_processor = FileProcessor()


import os
import logging
import mimetypes
from typing import Tuple

# ---------- Optional python-magic (libmagic is often missing on Azure) ----------
try:
    import magic  # type: ignore
    HAS_MAGIC = True
except Exception:
    magic = None  # type: ignore
    HAS_MAGIC = False

import pandas as pd
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from docx import Document as DocxDocument  # NEW: robust DOCX reader (pip install python-docx)

# ---------- LangChain loaders (v0.3+ first, fallback to older paths) ----------
try:
    from langchain_community.document_loaders import (
        PyPDFLoader,
        TextLoader,
        UnstructuredWordDocumentLoader,
    )
except Exception:  # fallback for older langchain installs
    from langchain.document_loaders import PyPDFLoader, TextLoader  # type: ignore
    from langchain.document_loaders.word_document import (  # type: ignore
        UnstructuredWordDocumentLoader,
    )

from langchain.text_splitter import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)


# ---- Robust fallback by file extension (when magic/mimetypes are unreliable) ----
def guess_type_by_extension(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()

    # Common office/doc types
    if ext == ".docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if ext == ".doc":
        return "application/msword"
    if ext == ".pdf":
        return "application/pdf"
    if ext == ".txt":
        return "text/plain"
    if ext == ".csv":
        return "text/csv"
    if ext == ".xlsx":
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if ext == ".xls":
        return "application/vnd.ms-excel"
    if ext == ".pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if ext == ".ppt":
        return "application/vnd.ms-powerpoint"

    # Common images
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"):
        return f"image/{ext.lstrip('.')}"

    return "application/octet-stream"


# Normalize office MIME even if 'magic' returns application/zip
def normalize_office_mime(filename: str, file_type: str | None) -> str:
    ext = os.path.splitext(filename)[1].lower()
    # When magic reports ZIP or unknown, prefer OOXML based on extension
    if file_type in (None, "", "application/octet-stream", "application/zip"):
        if ext == ".docx":
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if ext == ".xlsx":
            return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        if ext == ".pptx":
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if ext == ".doc":
            return "application/msword"
        if ext == ".xls":
            return "application/vnd.ms-excel"
        if ext == ".ppt":
            return "application/vnd.ms-powerpoint"

    # Prefer new OOXML if extension says so but mimetype is legacy
    if file_type == "application/msword" and ext == ".docx":
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if file_type == "application/vnd.ms-excel" and ext == ".xlsx":
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if file_type == "application/vnd.ms-powerpoint" and ext == ".pptx":
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    return file_type or "application/octet-stream"


class FileProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    async def process_file(self, file_path: str, filename: str) -> Tuple[str, str]:
        """Process uploaded file and extract content summary"""
        try:
            # Get candidates from both detectors
            mt_guess, _ = mimetypes.guess_type(filename)
            mg_guess = None
            if HAS_MAGIC:
                try:
                    mg_guess = magic.from_file(file_path, mime=True)  # type: ignore
                except Exception as e:
                    logger.warning(f"[magic] failed to detect file type: {e}")

            # Start with extension map (never 'zip' for OOXML)
            file_type = guess_type_by_extension(filename)

            # If magic confidently says image/* or pdf/text, prefer it
            if mg_guess and (mg_guess.startswith("image/") or mg_guess in ("application/pdf",) or mg_guess.startswith("text/")):
                file_type = mg_guess

            # If still generic from extension, let mimetypes refine
            if file_type == "application/octet-stream" and mt_guess:
                file_type = mt_guess

            # Finally, normalize office doc zip/legacy types by extension
            file_type = normalize_office_mime(filename, file_type)

            logger.info(f"[MIME] filename={filename} ext_guess={guess_type_by_extension(filename)} mt_guess={mt_guess} magic={mg_guess} -> final={file_type}")

            # ---- Dispatch by final file_type ----
            if file_type.startswith("image/"):
                return await self._process_image(file_path, filename)
            elif file_type == "application/pdf":
                return await self._process_pdf(file_path, filename)
            elif file_type in (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "application/msword",
            ):
                return await self._process_word(file_path, filename)
            elif file_type == "text/plain":
                return await self._process_text(file_path, filename)
            elif file_type in (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
            ):
                return await self._process_excel(file_path, filename)
            elif file_type in (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
            ):
                return await self._process_ppt(file_path, filename)
            elif file_type == "application/octet-stream":
                return await self._process_unknown(file_path, filename, file_type)
            else:
                return (f"Unsupported file type: {file_type}", "error")

        except Exception as e:
            logger.exception(f"Error processing file {filename}: {e}")
            return (f"Error processing file: {str(e)}", "error")

    # ---------------- Handlers ----------------

    async def _process_image(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            with Image.open(file_path) as img:
                img.verify()
            with Image.open(file_path) as img2:
                width, height = img2.size
                format_name = img2.format

            content = (
                f"Image uploaded: {filename}\n"
                f"Format: {format_name}\n"
                f"Dimensions: {width}x{height} pixels\n"
                "Ready for analysis. You can ask me to describe the image or answer questions about it."
            )
            return content, "image"
        except UnidentifiedImageError:
            return "Error processing image: unrecognized or corrupt image format.", "error"
        except Exception as e:
            return f"Error processing image: {str(e)}", "error"

    async def _process_pdf(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
            if not documents:
                return "PDF appears to be empty or unreadable.", "error"

            full_text = "\n\n".join(doc.page_content for doc in documents)
            preview = full_text[:2000] + ("..." if len(full_text) > 2000 else "")
            content = (
                f"PDF Document: {filename}\n"
                f"Pages: {len(documents)}\n"
                f"Content length: {len(full_text)} characters\n\n"
                "Document preview:\n" + preview + "\n\n"
                "Document loaded successfully. You can ask me questions about its content."
            )
            return content, "pdf"
        except Exception as e:
            return f"Error processing PDF: {str(e)}", "error"

    async def _process_word(self, file_path: str, filename: str) -> Tuple[str, str]:
        # Prefer python-docx first (no heavy deps, very reliable for .docx)
        try:
            doc = DocxDocument(file_path)
            paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
            full_text = "\n".join(paras)

            if not full_text:
                # try tables if body text is empty
                table_lines = []
                for tbl in doc.tables:
                    for row in tbl.rows:
                        cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                        if cells:
                            table_lines.append(" | ".join(cells))
                full_text = "\n".join(table_lines)

            if full_text:
                preview = full_text[:2000] + ("..." if len(full_text) > 2000 else "")
                content = (
                    f"Word Document: {filename}\n"
                    f"Content length: {len(full_text)} characters\n\n"
                    "Document preview:\n" + preview + "\n\n"
                    "Document loaded successfully. You can ask me questions about its content."
                )
                return content, "word"
        except Exception as e:
            logger.warning(f"[python-docx] failed on {filename}: {e}")

        # Fallback to Unstructured loader if python-docx failed
        try:
            loader = UnstructuredWordDocumentLoader(file_path)
            documents = loader.load()
            if not documents:
                return "Word document appears to be empty or unreadable.", "error"

            full_text = "\n\n".join(doc.page_content for doc in documents)
            preview = full_text[:2000] + ("..." if len(full_text) > 2000 else "")
            content = (
                f"Word Document: {filename}\n"
                f"Content length: {len(full_text)} characters\n\n"
                "Document preview:\n" + preview + "\n\n"
                "Document loaded successfully. You can ask me questions about its content."
            )
            return content, "word"
        except Exception as e:
            return f"Error processing Word document: {str(e)}", "error"

    async def _process_text(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            try:
                loader = TextLoader(file_path, encoding="utf-8", autodetect_encoding=True, errors="ignore")
            except TypeError:
                loader = TextLoader(file_path, encoding="utf-8")
            documents = loader.load()
            if not documents:
                return "Text file appears to be empty.", "error"

            full_text = documents[0].page_content
            preview = full_text[:2000] + ("..." if len(full_text) > 2000 else "")
            content = (
                f"Text File: {filename}\n"
                f"Content length: {len(full_text)} characters\n\n"
                "File preview:\n" + preview + "\n\n"
                "File loaded successfully. You can ask me questions about its content."
            )
            return content, "text"
        except Exception as e:
            return f"Error processing text file: {str(e)}", "error"

    async def _process_csv(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            df = pd.read_csv(file_path, nrows=5, encoding="utf-8", on_bad_lines="skip")
            preview = df.to_string(index=False)
            content = (
                f"CSV File: {filename}\n"
                f"Preview (first 5 rows):\n{preview}\n\n"
                "CSV file loaded successfully. You can ask me about its contents."
            )
            return content, "csv"
        except UnicodeDecodeError:
            try:
                df = pd.read_csv(file_path, nrows=5, encoding="latin-1", on_bad_lines="skip")
                preview = df.to_string(index=False)
                content = (
                    f"CSV File: {filename}\n"
                    f"(latin-1 fallback)\n"
                    f"Preview (first 5 rows):\n{preview}\n\n"
                    "CSV file loaded successfully. You can ask me about its contents."
                )
                return content, "csv"
            except Exception as e:
                return f"Error processing CSV file (encoding): {str(e)}", "error"
        except Exception as e:
            return f"Error processing CSV file: {str(e)}", "error"

    async def _process_excel(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            engine = None
            try:
                import openpyxl  # noqa: F401
                engine = "openpyxl"
            except Exception:
                pass

            df = pd.read_excel(file_path, engine=engine)
            preview = df.head(5).to_string(index=False)
            content = (
                f"Excel File: {filename}\n"
                f"Rows: {df.shape[0]}, Columns: {df.shape[1]}\n\n"
                f"Preview:\n{preview}\n\n"
                "Excel file loaded successfully. You can ask me about its contents."
            )
            return content, "excel"
        except ImportError as e:
            return f"Error processing Excel file: {str(e)} (try installing 'openpyxl')", "error"
        except Exception as e:
            return f"Error processing Excel file: {str(e)}", "error"

    async def _process_ppt(self, file_path: str, filename: str) -> Tuple[str, str]:
        try:
            prs = Presentation(file_path)  # note: supports .pptx, not legacy .ppt
            text_lines = []
            for i, slide in enumerate(prs.slides[:3]):
                lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        lines.append(shape.text.strip())
                if lines:
                    text_lines.append(f"Slide {i+1}: " + " | ".join(lines))
            flat = "\n".join(text_lines)
            preview = flat[:2000] + ("...\n\n" if len(flat) > 2000 else "\n")
            content = (
                f"PowerPoint File: {filename}\n"
                f"Slides: {len(prs.slides)}\n\n"
                f"{preview}"
                "PowerPoint file loaded. Ask me about its content."
            )
            return content, "ppt"
        except Exception as e:
            # If user uploaded legacy .ppt, suggest converting
            if os.path.splitext(filename)[1].lower() == ".ppt":
                return ("This looks like a legacy .ppt file. Please convert to .pptx and re-upload.", "error")
            return f"Error processing PowerPoint file: {str(e)}", "error"

    async def _process_unknown(self, file_path: str, filename: str, file_type: str) -> Tuple[str, str]:
        try:
            size_kb = round(os.path.getsize(file_path) / 1024, 2)
            content = (
                f"File uploaded: {filename}\n"
                f"Type: {file_type}\n"
                f"Size: {size_kb} KB\n\n"
                "This file type is not supported for detailed parsing. You can try converting it to PDF, TXT, or DOCX."
            )
            return content, "unsupported"
        except Exception as e:
            return f"Error analyzing unknown file type: {str(e)}", "error"


# Export instance
file_processor = FileProcessor()
